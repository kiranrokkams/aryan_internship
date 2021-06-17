from ..contents.Component import Component
from . import layouts
from pptx.slide import SlideLayout
from .contents.misc import PptText
from .contents.PptTable import PptTable
from .contents.PptImage import PptImage
from .contents.PptChart import PptChart
import pandas as pd
from pandas.io.formats.style import Styler
import tigerml
from tigerml.core.dataframe.helpers import detigerify

def wrap_content(content, content_name=''):
	if isinstance(content, tigerml.core.dataframe.dataframe.DataFrame):
		content = detigerify(content)
	if isinstance(content, str):
		content = PptText(content)
	elif content.__module__.startswith('tigerml.core.reports.ppt.contents'):
		content = content
	elif isinstance(content, pd.DataFrame) or isinstance(content, Styler):
		content = PptTable(content, title=content_name)
	else:
		content = PptImage(content, name=content_name)
	return content


class SlideComponent(Component):

	def __init__(self, content, slide):
		# import pdb
		# pdb.set_trace()
		# assert content.__module__.startswith('tigerml.core.reports.ppt') or content.__module__.startswith('pptx')
		content = wrap_content(content)
		super().__init__(content, parent=slide)
		self.slide = slide

	def save(self, placeholder, slide_obj):
		self.content.save(placeholder, slide_obj)


class Slide:

	def __init__(self, title='', layout=None):
		if not layout:
			layout = layouts.title_and_content
		self.set_layout(layout)
		self.title = title
		self.contents = list()

	@property
	def pretty_title(self):
		from tigerml.core.utils import prettify_slug
		return prettify_slug(self.title)

	def set_layout(self, layout):
		assert isinstance(layout, SlideLayout), 'Unsopported layout for slide. Should be one of {}'.format(layouts.list)
		self.layout = layout
		return self

	def add_content(self, content):
		self.contents.append(SlideComponent(content, self))

	def save(self, ppt, parent_name=None):
		if parent_name:
			section_name = parent_name + ((' - ' + self.pretty_title) if self.pretty_title else '')
		else:
			section_name = self.pretty_title
		slide_obj = ppt.slides.add_slide(slide_layout=self.layout)
		if self.title and slide_obj.shapes.title:
			slide_obj.shapes.title.text = section_name
		for index, component in enumerate(self.contents):
			if index + 1 < len(slide_obj.placeholders):
				placeholder = slide_obj.placeholders[index+1]
				component.save(placeholder, slide_obj)
			else:
				raise Exception('Not enough shapes in the slide for the contents.')
		return slide_obj

