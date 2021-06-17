from tigerml.core.utils import check_or_create_path
from .ContentGroup import ContentGroup
from .Slide import Slide
from . import layouts
import pptx


class PptReport(ContentGroup):

	def __init__(self, title, filename=''):
		if filename:
			assert '.pptx' in filename, 'name should have the extension .pptx'
		super().__init__()
		self._filename = filename
		self.title = title
		self._title_slide = Slide(title=self.title, layout=layouts.title_slide)

	@property
	def filename(self):
		return (self._filename + ('.pptx' if self._filename[:-5] != '.pptx' else '')) if self._filename else \
		       (self.title + ('.pptx' if self.title[:-5] != '.pptx' else ''))

	@property
	def slides(self):
		all_slides = super().slides
		if self.title:
			all_slides = [self._title_slide] + all_slides
		return all_slides

	def save(self, path='', tiger_template=False):
		template_path = None
		if tiger_template:
			import os
			_thisdir = os.path.split(__file__)[0]
			template_path = os.path.join(_thisdir, "tiger_template.pptx")
		ppt = pptx.Presentation(template_path)
		if path:
			if path[-1] != '/':
				path = path + '/'
			check_or_create_path(path)
		if self.title:
			self._title_slide.save(ppt)
		for content in self.contents:
			content.save(ppt)
		ppt.save(path + self.filename)


class Section(ContentGroup):

	def __init__(self, title=''):
		super().__init__()
		self._title = title
		self.show_title = True
		self._title_slide = Slide(title=self.pretty_title, layout=layouts.section_header)

	def set_title(self, title):
		self._title = title
		self._title_slide.title = self.pretty_title
		return self

	@property
	def pretty_title(self):
		from tigerml.core.utils import prettify_slug
		return prettify_slug(self.title)

	@property
	def title(self):
		return self._title

	@property
	def title_slide(self):
		# if title:
		# 	self._title_slide.title = title
		return self._title_slide

	@property
	def slides(self):
		all_slides = super().slides
		if self.title:
			all_slides = [self._title_slide] + all_slides
		return all_slides

	def save(self, ppt, parent_name=''):
		if parent_name:
			section_name = parent_name + ((' - ' + self.pretty_title) if self.pretty_title else '')
		else:
			section_name = self.pretty_title
		s_objs = list()
		if self.title and self.show_title:
			self.title_slide.save(ppt, parent_name=parent_name)
		for content in self.contents:
			s_obj = content.save(ppt, parent_name=self.pretty_title or section_name)
			s_objs.append(s_obj)
		return s_objs

